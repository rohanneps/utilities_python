import argparse
import urllib.request
import os
from pptx import Presentation 
from pptx.util import Inches
import os
from bs4 import BeautifulSoup
import shutil

work_dir = os.path.dirname(__file__)

IMG_FOLDER = 'ppt_images'

if not os.path.exists(IMG_FOLDER):
	os.makedirs(IMG_FOLDER)
else:
	os.system('rm -r {}'.format(IMG_FOLDER))
	os.makedirs(IMG_FOLDER)

def downloadImage(image_url, image_name):
	f = open(image_name,'wb')
	# imageContent = requests.get(url).content
	imageContent = urllib.request.urlopen(image_url).read()
	f.write(imageContent)
	f.close()


def downloadPptImages(url):
	html = urllib.request.urlopen(url).read()
	soup = BeautifulSoup(html, 'lxml')
	title = IMG_FOLDER  # soup.title.string
	images = soup.findAll('img', {'class': 'slide_image'})

	c = 1
	for image in images:
		image_url = image.get('data-full').split('?')[0]
		name = '{}.jpg'.format(c)
		# for ubuntu
		# command = 'wget %s -O %s --quiet' % (image_url,name)
		# os.system(command)

		downloadImage(image_url, name)
		shutil.move(name,os.path.join(IMG_FOLDER,name))
		c += 1


if __name__ == "__main__":
	parser = argparse.ArgumentParser()
	parser.add_argument("url", type=str,
	                    help="download an slideshare presentation given the url")
	args = parser.parse_args()

	downloadPptImages(args.url)

	prs = Presentation() 


	for (dirpath, dirnames, filenames) in os.walk(IMG_FOLDER):
		total_file_count = (len(filenames))
	print('Number of slides:{}'.format(total_file_count))

	for i in range(total_file_count):
		filename = '{}.jpg'.format(i+1)
		img_path = os.path.join(IMG_FOLDER,filename)
		blank_slide_layout = prs.slide_layouts[6] 
		slide = prs.slides.add_slide(blank_slide_layout)

		# left = top = Inches(1)
		pic_left  = int(prs.slide_width * 0.12)
		pic_top   = int(prs.slide_height * 0.12)
		pic_width = int(prs.slide_width * 0.7)
		pic_height = int(prs.slide_height * 0.7)
		pic = slide.shapes.add_picture(img_path, pic_left, pic_top, pic_width, pic_height) 

	prs.save('presentation.pptx')
	os.system('rm -r {}'.format(IMG_FOLDER))